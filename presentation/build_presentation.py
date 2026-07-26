#!/usr/bin/env python3
"""
Build the official AdOS executive presentation.

Outputs (in this directory):
  - AdOS_Executive_Presentation.pptx  (editable, vector shapes, 16:9, dark theme)
  - AdOS_Executive_Presentation.pdf   (exported, same content)

Design: brand dark theme, one accent per slide, native vector shapes (the
perimeter / pipeline / comparison diagrams are drawn with shapes, not images),
bilingual EN/TR. No external or copyrighted assets. Animation-ready: each build
element is a discrete shape so PowerPoint animations can be added per object.

Run:  python3 build_presentation.py
Deps: python-pptx, reportlab (both already available in the build environment).
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX_OUT = os.path.join(HERE, "AdOS_Executive_Presentation.pptx")
PDF_OUT = os.path.join(HERE, "AdOS_Executive_Presentation.pdf")

# Brand palette (hex ints)
INK = RGBColor(0x0E, 0x11, 0x16)
PANEL = RGBColor(0x16, 0x1B, 0x22)
PANEL2 = RGBColor(0x1C, 0x22, 0x30)
LINE = RGBColor(0x2A, 0x31, 0x40)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x98, 0xA9)
BRAND = RGBColor(0x5B, 0x8C, 0xFF)
VIOLET = RGBColor(0x9D, 0x7B, 0xFF)
OK = RGBColor(0x3F, 0xB9, 0x50)
FONT = "Arial"

# ── Slide content model ──────────────────────────────────────────────────────
# kind: title | content | diagram | security | audience | compare | demo | closing
SLIDES = [
    {"kind": "title", "n": 1,
     "en_t": "AdOS", "tr_t": "AdOS",
     "en_s": "The Advertising Operating System", "tr_s": "Reklam İşletim Sistemi",
     "en_l": "Enterprise AI that never leaves your building.",
     "tr_l": "Binanızdan hiç çıkmayan kurumsal yapay zekâ."},

    {"kind": "content", "n": 2, "sec": "The world changed",
     "en_t": "The problems every organization faces", "tr_t": "Her kurumun karşılaştığı sorunlar",
     "en_l": "Work is slow. Knowledge is scattered. Approvals wait on busy people. And AI promises help — at the price of your data.",
     "tr_l": "İşler yavaş. Bilgi dağınık. Onaylar meşgul kişileri bekliyor. Yapay zekâ yardım vaat ediyor — verilerinizi karşılığında isteyerek."},

    {"kind": "content", "n": 3, "sec": "The world changed",
     "en_t": "Digital transformation is changing shape", "tr_t": "Dijital dönüşüm biçim değiştiriyor",
     "en_l": "You digitized your processes. The next step is intelligence — systems that understand your business and do the work, on your terms.",
     "tr_l": "Süreçlerinizi dijitalleştirdiniz. Sıradaki adım zekâ — işinizi anlayan ve işi yapan sistemler, sizin koşullarınızda."},

    {"kind": "content", "n": 4, "sec": "The world changed",
     "en_t": "Your most valuable asset walks out daily", "tr_t": "En değerli varlığınız her gün çıkıyor",
     "en_l": "Experience retires. Documents get lost. Every departure takes knowledge with it. What if your organization could remember everything it learned?",
     "tr_l": "Deneyim emekli oluyor. Belgeler kayboluyor. Her ayrılış bilgiyi götürüyor. Ya kurumunuz öğrendiği her şeyi hatırlayabilseydi?"},

    {"kind": "content", "n": 5, "sec": "The world changed",
     "en_t": "AI is inevitable. The risk is where it runs.", "tr_t": "Yapay zekâ kaçınılmaz. Risk, nerede çalıştığında.",
     "en_l": "The question is no longer whether to adopt AI — it's whether you adopt it safely, with your data under your control.",
     "tr_l": "Soru artık yapay zekâyı benimseyip benimsememek değil — onu güvenle, verileriniz kontrolünüzde benimseyip benimsemediğiniz."},

    {"kind": "content", "n": 6, "sec": "The answer",
     "en_t": "Local AI: the intelligence comes to your data", "tr_t": "Yerel yapay zekâ: zekâ verinize gelir",
     "en_l": "The AI models run on your own computers. No cloud. No API keys. Your information never goes to anyone else.",
     "tr_l": "Yapay zekâ modelleri kendi bilgisayarlarınızda çalışır. Bulut yok. API anahtarı yok. Bilgileriniz asla başkasına gitmez."},

    {"kind": "diagram", "diagram": "perimeter", "n": 7, "sec": "The answer",
     "en_t": "One boundary: everything stays inside", "tr_t": "Tek sınır: her şey içeride kalır",
     "en_l": "A request goes to a local AI manager, which talks only to a model on your hardware. Nothing crosses the line.",
     "tr_l": "İstek yerel bir yapay zekâ yöneticisine gider; o da yalnızca donanımınızdaki bir modelle konuşur. Hiçbir şey sınırı geçmez."},

    {"kind": "diagram", "diagram": "sovereignty", "n": 8, "sec": "The answer",
     "en_t": "Your data never leaves your building", "tr_t": "Verileriniz binanızdan hiç çıkmaz",
     "en_l": "There is nowhere for it to leak, because nothing goes out.",
     "tr_l": "Sızabileceği bir yer yoktur, çünkü dışarı hiçbir şey çıkmaz."},

    {"kind": "security", "n": 9, "sec": "The answer",
     "en_t": "Security by architecture, not by promise", "tr_t": "Vaatle değil, mimariyle güvenlik",
     "cards": [
         ("Separation", "Ayrım", "Each unit's data is walled off.", "Her birimin verisi ayrılmıştır."),
         ("Protection", "Koruma", "Everything is encrypted and recoverable.", "Her şey şifreli ve kurtarılabilir."),
         ("Accountability", "Hesap verebilirlik", "Every action is recorded.", "Her eylem kaydedilir."),
     ]},

    {"kind": "diagram", "diagram": "pipeline", "n": 10, "sec": "The answer",
     "en_t": "The work gets done — you approve each step", "tr_t": "İş yapılır — her adımı siz onaylarsınız",
     "en_l": "AdOS moves work through clear stages, with a human decision at each gate. Autonomous, but never unaccountable.",
     "tr_l": "AdOS işi net aşamalardan geçirir; her adımda bir insan kararı vardır. Otonom, ama asla hesapsız değil."},

    {"kind": "content", "n": 11, "sec": "The answer",
     "en_t": "The Company Brain: your living memory", "tr_t": "Şirket Beyni: yaşayan hafızanız",
     "en_l": "Every document, decision and outcome becomes shared, searchable knowledge that grows more valuable over time — and never leaves your walls.",
     "tr_l": "Her belge, karar ve sonuç; zamanla değerlenen, paylaşılan ve aranabilir bir bilgiye dönüşür — ve asla duvarlarınızdan çıkmaz."},

    {"kind": "content", "n": 12, "sec": "The answer",
     "en_t": "Digital Employees for every department", "tr_t": "Her departman için Dijital Çalışanlar",
     "en_l": "Specialized assistants — HR, Finance, Quality, Maintenance and more — that know your policies and help your people, within clear limits.",
     "tr_l": "Uzmanlaşmış asistanlar — İK, Finans, Kalite, Bakım ve daha fazlası — politikalarınızı bilen ve net sınırlar içinde yardım eden."},

    {"kind": "content", "n": 13, "sec": "Believe it",
     "en_t": "What it does, in the real world", "tr_t": "Gerçek dünyada ne yapar",
     "en_l": "Run advertising campaigns. Answer policy questions instantly. Route approvals. Draft reports. Find the right document in seconds.",
     "tr_l": "Reklam kampanyaları yürütün. Politika sorularını anında yanıtlayın. Onayları yönlendirin. Rapor hazırlayın. Doğru belgeyi saniyeler içinde bulun."},

    {"kind": "audience", "n": 14, "sec": "Believe it",
     "en_t": "Built for organizations that can't compromise", "tr_t": "Ödün veremeyen kurumlar için",
     "cards": [
         ("Enterprises", "Kurumlar"),
         ("Public institutions", "Kamu kurumları"),
         ("Industrial zones (OIZ)", "Organize sanayi bölgeleri"),
         ("Agencies", "Ajanslar"),
     ]},

    {"kind": "content", "n": 15, "sec": "Believe it",
     "en_t": "The return: more output, predictable cost, no data risk", "tr_t": "Getiri: daha fazla çıktı, öngörülebilir maliyet, veri riski yok",
     "en_l": "Your team produces more with the same headcount. AI costs are fixed — no per-use fees. And you take on zero data-exposure risk.",
     "tr_l": "Ekibiniz aynı personelle daha fazla üretir. Yapay zekâ maliyetleri sabittir — kullanım başına ücret yok. Ve sıfır veri-ifşa riski alırsınız."},

    {"kind": "content", "n": 16, "sec": "Believe it",
     "en_t": "It runs where you already run software", "tr_t": "Zaten yazılım çalıştırdığınız yerde çalışır",
     "en_l": "Your data center, private cloud or on-site servers. We install and guide; you operate and own it. Open components, no lock-in.",
     "tr_l": "Veri merkeziniz, özel bulutunuz ya da yerinde sunucularınız. Biz kurar ve yol gösteririz; siz işletir ve sahiplenirsiniz. Açık bileşenler, bağımlılık yok."},

    {"kind": "demo", "n": 17, "sec": "Believe it",
     "en_t": "See it work — locally", "tr_t": "Çalışırken görün — yerelde",
     "en_l": "A real campaign, produced entirely on local infrastructure, with no internet connection. Nothing you see leaves the building.",
     "tr_l": "Tamamen yerel altyapıda, internet bağlantısı olmadan üretilen gerçek bir kampanya. Gördüğünüz hiçbir şey binadan çıkmaz."},

    {"kind": "compare", "n": 18, "sec": "Believe it",
     "en_t": "Everyone can generate output. Only AdOS keeps your data.", "tr_t": "Herkes çıktı üretebilir. Yalnızca AdOS verinizi korur.",
     "rows": [
         ("Your data", "Verileriniz", "Leaves to the cloud", "Buluta çıkar", "Stays inside", "İçeride kalır"),
         ("Cost", "Maliyet", "Per-use fees", "Kullanım başına ücret", "Fixed, predictable", "Sabit, öngörülebilir"),
         ("Control", "Kontrol", "Vendor's terms", "Tedarikçinin koşulları", "You own it", "Size ait"),
         ("Compliance", "Uyum", "Exposure", "Risk altında", "By design", "Tasarımdan gelir"),
     ]},

    {"kind": "content", "n": 19, "sec": "Decide",
     "en_t": "A staged, low-risk path", "tr_t": "Aşamalı, düşük riskli bir yol",
     "en_l": "Pilot on your infrastructure → prove value on your data → roll out → compounding advantage as your knowledge grows.",
     "tr_l": "Altyapınızda pilot → verinizde değeri kanıtla → yaygınlaştır → bilginiz büyüdükçe bileşik avantaj."},

    {"kind": "content", "n": 20, "sec": "Decide",
     "en_t": "In one picture", "tr_t": "Tek bir resimde",
     "en_l": "Modern AI for advertising and enterprise knowledge, running entirely on your infrastructure, under your control, with your data never leaving your building.",
     "tr_l": "Reklam ve kurumsal bilgi için modern yapay zekâ; tamamen altyapınızda, kontrolünüzde çalışan, verileriniz binanızdan hiç çıkmadan."},

    {"kind": "closing", "n": 21, "sec": "Decide",
     "en_t": "Let's run a pilot on your infrastructure", "tr_t": "Altyapınızda bir pilot yürütelim",
     "en_l": "A scoped pilot, on your own hardware, guided by us — prove every point with your own data staying inside.",
     "tr_l": "Kapsamı belirli bir pilot, kendi donanımınızda, bizim rehberliğimizde — her noktayı, kendi veriniz içeride kalarak kanıtlayın."},

    {"kind": "title", "n": 22,
     "en_t": "Thank you", "tr_t": "Teşekkürler",
     "en_s": "AdOS — The Advertising Operating System", "tr_s": "AdOS — Reklam İşletim Sistemi",
     "en_l": "The AI comes to your data — never the other way around.",
     "tr_l": "Yapay zekâ verinize gelir — asla tersi değil."},
]

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ── PPTX helpers ─────────────────────────────────────────────────────────────
def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
    return tb


def _rect(slide, shape, left, top, width, height, fill, line_color, line_w=1.5):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _node(slide, left, top, width, height, en, tr, line_color=LINE, fill=PANEL):
    sp = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
               fill, line_color)
    tf = sp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = en; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = TEXT; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = tr; r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED; r2.font.name = FONT
    return sp


def _connect(slide, x1, y1, x2, y2, color=BRAND):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(1.75)
    return cn


def _chrome(slide, s):
    """Accent bar, section eyebrow, slide number, footer brand."""
    _rect(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55),
          Inches(0.5), Inches(0.09), BRAND, None)
    if s.get("sec"):
        _text(slide, Inches(0.6), Inches(0.72), Inches(8), Inches(0.4),
              [(s["sec"].upper(), 11, BRAND, True)])
    _text(slide, Inches(12.0), Inches(6.95), Inches(1.1), Inches(0.4),
          [(f"{s['n']} / {len(SLIDES)}", 10, MUTED, False)], align=PP_ALIGN.RIGHT)
    _text(slide, Inches(0.6), Inches(6.95), Inches(4), Inches(0.4),
          [("AdOS — The Advertising Operating System", 10, MUTED, False)])


def _titles(slide, s, top=Inches(1.4)):
    _text(slide, Inches(0.6), top, Inches(12.1), Inches(1.6),
          [(s["en_t"], 34, TEXT, True), (s["tr_t"], 20, MUTED, False)])


def _perimeter(slide, top=Inches(3.2), big_number=False):
    bx, by, bw, bh = Inches(1.2), top, Inches(10.9), Inches(2.7)
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh, None, BRAND, 1.75)
    _text(slide, bx + Inches(0.2), by + Inches(0.08), Inches(4), Inches(0.4),
          [("YOUR NETWORK · AĞINIZ", 10, MUTED, True)])
    nodes = [("Your request", "İsteğiniz"),
             ("AI Manager (local)", "Yerel yönetici"),
             ("Local model", "Yerel model")]
    nw, nh = Inches(2.7), Inches(1.1)
    ny = by + Inches(1.1)
    xs = [bx + Inches(0.5), bx + Inches(4.1), bx + Inches(7.7)]
    for (en, tr), x in zip(nodes, xs):
        lc = VIOLET if en.startswith("AI") else LINE
        _node(slide, x, ny, nw, nh, en, tr, line_color=lc)
    midy = ny + nh / 2
    _connect(slide, xs[0] + nw, midy, xs[1], midy)
    _connect(slide, xs[1] + nw, midy, xs[2], midy)
    # boundary "nothing crosses" line
    ly = by + bh + Inches(0.15)
    _rect(slide, MSO_SHAPE.RECTANGLE, bx, ly, bw, Emu(9525), VIOLET, None)
    _text(slide, bx, ly + Inches(0.05), bw, Inches(0.35),
          [("Nothing crosses the line · Hiçbir şey sınırı geçmez", 11, VIOLET, True)],
          align=PP_ALIGN.CENTER)
    if big_number:
        _text(slide, Inches(0.6), Inches(2.95), Inches(12.1), Inches(0.6),
              [("0 data leaves your network  ·  Ağınızdan 0 veri çıkar", 24, BRAND, True)],
              align=PP_ALIGN.CENTER)


def _pipeline(slide, top=Inches(3.6)):
    stages = [("Brief", "Brif"), ("Creative", "Kreatif"), ("Campaign", "Kampanya"),
              ("Results", "Sonuçlar"), ("Executive", "Yönetim")]
    nw_in, gap_in = 2.15, 0.25
    nw, nh = Inches(nw_in), Inches(1.2)
    gap = Inches(gap_in)
    total_in = nw_in * 5 + gap_in * 4
    x = Inches((13.333 - total_in) / 2)
    midy = top + nh / 2
    prev_right = None
    for en, tr in stages:
        _node(slide, x, top, nw, nh, en, tr, line_color=LINE)
        if prev_right is not None:
            # green approval gate between stages
            gx = prev_right + (x - prev_right) / 2
            g = _rect(slide, MSO_SHAPE.OVAL, gx - Inches(0.16), midy - Inches(0.16),
                      Inches(0.32), Inches(0.32), INK, OK, 1.5)
            tf = g.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "✓"; r.font.size = Pt(12); r.font.bold = True
            r.font.color.rgb = OK; r.font.name = "Calibri"
        prev_right = x + nw
        x = x + nw + gap
    _text(slide, Inches(0.6), top + nh + Inches(0.45), Inches(12.1), Inches(0.5),
          [("You approve every stage before the next begins · Bir sonraki adım başlamadan her aşamayı siz onaylarsınız", 13, MUTED, False)],
          align=PP_ALIGN.CENTER)


def _cards(slide, cards, top=Inches(3.4)):
    n = len(cards)
    cw_in, gap_in = (2.7 if n == 4 else 3.6), 0.4
    cw = Inches(cw_in)
    gap = Inches(gap_in)
    total_in = cw_in * n + gap_in * (n - 1)
    x = Inches((13.333 - total_in) / 2)
    for card in cards:
        sp = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, cw, Inches(2.4),
                   PANEL, LINE)
        tf = sp.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = card[0]
        r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = TEXT
        r.font.name = FONT
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = card[1]
        r2.font.size = Pt(12); r2.font.color.rgb = BRAND; r2.font.name = FONT
        if len(card) >= 4:
            p3 = tf.add_paragraph(); p3.space_before = Pt(8)
            r3 = p3.add_run(); r3.text = card[2]; r3.font.size = Pt(12)
            r3.font.color.rgb = MUTED; r3.font.name = FONT
            p4 = tf.add_paragraph(); r4 = p4.add_run(); r4.text = card[3]
            r4.font.size = Pt(11); r4.font.color.rgb = MUTED; r4.font.name = FONT
        x = x + cw + gap


def _compare(slide, rows, top=Inches(3.0)):
    lx, rx = Inches(1.2), Inches(7.1)
    cw = Inches(5.0)
    _text(slide, lx, top - Inches(0.5), cw, Inches(0.4),
          [("Cloud AI tools · Bulut yapay zekâ araçları", 13, MUTED, True)])
    _text(slide, rx, top - Inches(0.5), cw, Inches(0.4),
          [("AdOS", 13, BRAND, True)])
    y = top
    rh = Inches(0.8)
    for label_en, label_tr, l_en, l_tr, r_en, r_tr in rows:
        _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, lx, y, cw, rh, PANEL, LINE)
        _text(slide, lx + Inches(0.2), y + Inches(0.1), cw - Inches(0.4), rh,
              [(f"{label_en}: {l_en}", 12, TEXT, False), (l_tr, 10, MUTED, False)])
        _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, cw, rh, PANEL2, BRAND)
        _text(slide, rx + Inches(0.2), y + Inches(0.1), cw - Inches(0.4), rh,
              [(f"{label_en}: {r_en}", 12, TEXT, True), (r_tr, 10, BRAND, False)])
        y = y + rh + Inches(0.15)


def build_pptx():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, INK)
        kind = s["kind"]

        if kind == "title":
            _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(2.3),
                  Inches(1.5), Inches(1.5), None, None)
            # brand mark
            mk = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.06), Inches(2.35),
                       Inches(1.2), Inches(1.2), BRAND, None)
            tf = mk.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "▲"; r.font.size = Pt(30); r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = FONT
            _text(slide, Inches(0.6), Inches(3.75), Inches(12.1), Inches(1.0),
                  [(s["en_t"], 54, TEXT, True)], align=PP_ALIGN.CENTER)
            _text(slide, Inches(0.6), Inches(4.75), Inches(12.1), Inches(0.6),
                  [(s["en_s"], 20, BRAND, True), (s["tr_s"], 16, MUTED, False)],
                  align=PP_ALIGN.CENTER)
            _text(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0),
                  [(s["en_l"], 16, TEXT, False), (s["tr_l"], 14, MUTED, False)],
                  align=PP_ALIGN.CENTER)
            _text(slide, Inches(0.6), Inches(6.95), Inches(4), Inches(0.4),
                  [("100% local · No cloud · No API keys", 10, MUTED, False)])
            continue

        _chrome(slide, s)
        _titles(slide, s)

        if kind == "content":
            _text(slide, Inches(0.6), Inches(3.3), Inches(11.6), Inches(2.6),
                  [(s["en_l"], 20, TEXT, False), ("", 8, TEXT, False),
                   (s["tr_l"], 16, MUTED, False)])
        elif kind == "diagram" and s["diagram"] == "perimeter":
            _text(slide, Inches(0.6), Inches(2.85), Inches(11.6), Inches(0.5),
                  [(s["en_l"], 14, MUTED, False)])
            _perimeter(slide, top=Inches(3.35))
        elif kind == "diagram" and s["diagram"] == "sovereignty":
            _perimeter(slide, top=Inches(3.7), big_number=True)
        elif kind == "diagram" and s["diagram"] == "pipeline":
            _text(slide, Inches(0.6), Inches(2.85), Inches(11.6), Inches(0.7),
                  [(s["en_l"], 14, MUTED, False)])
            _pipeline(slide, top=Inches(3.7))
        elif kind == "security":
            _cards(slide, s["cards"], top=Inches(3.5))
        elif kind == "audience":
            _cards(slide, [(a, b, "", "") for a, b in s["cards"]], top=Inches(3.5))
            _text(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.5),
                  [("One platform can serve many organizations — each fully isolated. · Tek platform birçok kurumu hizmet verebilir — her biri tam izole.", 12, MUTED, False)],
                  align=PP_ALIGN.CENTER)
        elif kind == "compare":
            _compare(slide, s["rows"], top=Inches(3.0))
        elif kind == "demo":
            frame = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(3.0),
                          Inches(8.9), Inches(3.2), PANEL, BRAND)
            tf = frame.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "▶  LIVE / RECORDED DEMO"
            r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = BRAND
            r.font.name = FONT
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = s["en_l"]; r2.font.size = Pt(13)
            r2.font.color.rgb = TEXT; r2.font.name = FONT
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run(); r3.text = s["tr_l"]; r3.font.size = Pt(11)
            r3.font.color.rgb = MUTED; r3.font.name = FONT
        elif kind == "closing":
            _text(slide, Inches(0.6), Inches(3.2), Inches(12.1), Inches(1.2),
                  [(s["en_l"], 18, TEXT, False), (s["tr_l"], 14, MUTED, False)])
            _perimeter(slide, top=Inches(4.6))

    prs.save(PPTX_OUT)
    print("wrote", PPTX_OUT)


# ── PDF (reportlab) ──────────────────────────────────────────────────────────
def build_pdf():
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import HexColor
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    W, H = 960, 540
    reg = "TR"
    bold = "TRB"
    candidates = [
        ("/System/Library/Fonts/Supplemental/Arial.ttf",
         "/System/Library/Fonts/Supplemental/Arial Bold.ttf"),
        ("/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
         "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
    ]
    for rp, bp in candidates:
        if os.path.exists(rp):
            pdfmetrics.registerFont(TTFont(reg, rp))
            pdfmetrics.registerFont(TTFont(bold, bp if os.path.exists(bp) else rp))
            break

    ink = HexColor("#0E1116"); text = HexColor("#E6EDF3"); muted = HexColor("#8B98A9")
    brand = HexColor("#5B8CFF"); violet = HexColor("#9D7BFF"); ok = HexColor("#3FB950")
    panel = HexColor("#161B22"); line = HexColor("#2A3140")

    c = canvas.Canvas(PDF_OUT, pagesize=(W, H))

    def wrap(txt, font, size, maxw):
        words = txt.split()
        lines, cur = [], ""
        for w in words:
            t = (cur + " " + w).strip()
            if pdfmetrics.stringWidth(t, font, size) <= maxw:
                cur = t
            else:
                if cur:
                    lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
        return lines

    def para(x, y, txt, font, size, color, maxw, leading):
        c.setFont(font, size); c.setFillColor(color)
        for ln in wrap(txt, font, size, maxw):
            c.drawString(x, y, ln); y -= leading
        return y

    def perimeter(cy, big=False):
        bx, bw, bh = 90, 780, 190
        by = cy
        c.setStrokeColor(brand); c.setLineWidth(1.75); c.setFillColor(ink)
        c.roundRect(bx, by, bw, bh, 12, stroke=1, fill=0)
        c.setFont(reg, 9); c.setFillColor(muted)
        c.drawString(bx + 12, by + bh - 16, "YOUR NETWORK · AĞINIZ")
        nodes = [("Your request", "İsteğiniz"), ("AI Manager (local)", "Yerel yönetici"),
                 ("Local model", "Yerel model")]
        nw, nh = 200, 78
        ny = by + 55
        xs = [bx + 40, bx + 290, bx + 540]
        for (en, tr), nx in zip(nodes, xs):
            c.setFillColor(panel)
            c.setStrokeColor(violet if en.startswith("AI") else line)
            c.roundRect(nx, ny, nw, nh, 10, stroke=1, fill=1)
            c.setFillColor(text); c.setFont(bold, 12)
            c.drawCentredString(nx + nw / 2, ny + nh / 2 + 4, en)
            c.setFillColor(muted); c.setFont(reg, 9)
            c.drawCentredString(nx + nw / 2, ny + nh / 2 - 12, tr)
        c.setStrokeColor(brand); c.setLineWidth(1.75)
        midy = ny + nh / 2
        c.line(xs[0] + nw, midy, xs[1], midy)
        c.line(xs[1] + nw, midy, xs[2], midy)
        c.setStrokeColor(violet); c.setLineWidth(1.2)
        c.line(bx, by - 14, bx + bw, by - 14)
        c.setFillColor(violet); c.setFont(bold, 10)
        c.drawCentredString(W / 2, by - 28,
                            "Nothing crosses the line · Hiçbir şey sınırı geçmez")
        if big:
            c.setFillColor(brand); c.setFont(bold, 22)
            c.drawCentredString(W / 2, by + bh + 26,
                                "0 data leaves your network · Ağınızdan 0 veri çıkar")

    def pipeline(cy):
        stages = [("Brief", "Brif"), ("Creative", "Kreatif"), ("Campaign", "Kampanya"),
                  ("Results", "Sonuçlar"), ("Executive", "Yönetim")]
        nw, nh, gap = 150, 84, 20
        total = nw * 5 + gap * 4
        x = (W - total) / 2
        midy = cy + nh / 2
        prev = None
        for en, tr in stages:
            c.setFillColor(panel); c.setStrokeColor(line)
            c.roundRect(x, cy, nw, nh, 10, stroke=1, fill=1)
            c.setFillColor(text); c.setFont(bold, 12)
            c.drawCentredString(x + nw / 2, cy + nh / 2 + 4, en)
            c.setFillColor(muted); c.setFont(reg, 9)
            c.drawCentredString(x + nw / 2, cy + nh / 2 - 12, tr)
            if prev is not None:
                gx = prev + (x - prev) / 2
                c.setStrokeColor(ok); c.setFillColor(ink); c.setLineWidth(1.4)
                c.circle(gx, midy, 11, stroke=1, fill=1)
                # vector checkmark
                c.setStrokeColor(ok); c.setLineWidth(2)
                c.line(gx - 4, midy, gx - 1, midy - 4)
                c.line(gx - 1, midy - 4, gx + 5, midy + 5)
            prev = x + nw
            x += nw + gap
        c.setFillColor(muted); c.setFont(reg, 11)
        c.drawCentredString(W / 2, cy - 24,
                            "You approve every stage before the next begins · "
                            "Her aşamayı bir sonraki başlamadan onaylarsınız")

    def chrome(s):
        c.setFillColor(brand); c.rect(40, H - 44, 36, 6, stroke=0, fill=1)
        if s.get("sec"):
            c.setFillColor(brand); c.setFont(bold, 10)
            c.drawString(40, H - 66, s["sec"].upper())
        c.setFillColor(muted); c.setFont(reg, 9)
        c.drawRightString(W - 40, 24, f"{s['n']} / {len(SLIDES)}")
        c.drawString(40, 24, "AdOS — The Advertising Operating System")

    def titles(s, y=H - 130):
        c.setFillColor(text); c.setFont(bold, 26)
        yy = para(40, y, s["en_t"], bold, 26, text, W - 80, 30)
        c.setFont(reg, 15)
        para(40, yy - 2, s["tr_t"], reg, 15, muted, W - 80, 20)

    for s in SLIDES:
        c.setFillColor(ink); c.rect(0, 0, W, H, stroke=0, fill=1)
        kind = s["kind"]
        if kind == "title":
            c.setFillColor(brand)
            c.roundRect(W / 2 - 45, H - 200, 90, 90, 14, stroke=0, fill=1)
            c.setFillColor(HexColor("#FFFFFF")); c.setFont(bold, 40)
            c.drawCentredString(W / 2, H - 170, "▲")
            c.setFillColor(text); c.setFont(bold, 46)
            c.drawCentredString(W / 2, H - 250, s["en_t"])
            c.setFillColor(brand); c.setFont(bold, 17)
            c.drawCentredString(W / 2, H - 285, s["en_s"])
            c.setFillColor(muted); c.setFont(reg, 13)
            c.drawCentredString(W / 2, H - 307, s["tr_s"])
            c.setFillColor(text); c.setFont(reg, 14)
            c.drawCentredString(W / 2, H - 345, s["en_l"])
            c.setFillColor(muted); c.setFont(reg, 12)
            c.drawCentredString(W / 2, H - 366, s["tr_l"])
            c.showPage()
            continue

        chrome(s); titles(s)
        if kind == "content":
            yy = para(40, H - 210, s["en_l"], reg, 18, text, W - 80, 26)
            para(40, yy - 12, s["tr_l"], reg, 14, muted, W - 80, 20)
        elif kind == "diagram" and s["diagram"] == "perimeter":
            para(40, H - 200, s["en_l"], reg, 13, muted, W - 80, 18)
            perimeter(120)
        elif kind == "diagram" and s["diagram"] == "sovereignty":
            c.setFillColor(brand); c.setFont(bold, 22)
            c.drawCentredString(W / 2, H - 205,
                                "0 data leaves your network · Ağınızdan 0 veri çıkar")
            perimeter(95, big=False)
        elif kind == "diagram" and s["diagram"] == "pipeline":
            para(40, H - 200, s["en_l"], reg, 13, muted, W - 80, 18)
            pipeline(150)
        elif kind in ("security", "audience"):
            cards = s["cards"]
            n = len(cards); cw = 190 if n == 4 else 250; gap = 30
            total = cw * n + gap * (n - 1); x = (W - total) / 2
            for card in cards:
                c.setFillColor(panel); c.setStrokeColor(line)
                c.roundRect(x, 150, cw, 150, 12, stroke=1, fill=1)
                c.setFillColor(text); c.setFont(bold, 15)
                c.drawString(x + 16, 270, card[0])
                c.setFillColor(brand); c.setFont(reg, 11)
                c.drawString(x + 16, 252, card[1])
                if len(card) >= 4 and card[2]:
                    c.setFillColor(muted); c.setFont(reg, 10)
                    para(x + 16, 228, card[2], reg, 10, muted, cw - 32, 14)
                    para(x + 16, 200, card[3], reg, 9, muted, cw - 32, 12)
                x += cw + gap
        elif kind == "compare":
            lx, rx, cw, rh = 90, 500, 370, 58
            c.setFillColor(muted); c.setFont(bold, 11)
            c.drawString(lx, 320, "Cloud AI tools · Bulut yapay zekâ araçları")
            c.setFillColor(brand); c.drawString(rx, 320, "AdOS")
            y = 300 - rh
            for le, lt, l_en, l_tr, r_en, r_tr in s["rows"]:
                c.setFillColor(panel); c.setStrokeColor(line)
                c.roundRect(lx, y, cw, rh, 8, stroke=1, fill=1)
                c.setFillColor(text); c.setFont(reg, 11)
                c.drawString(lx + 12, y + rh - 22, f"{le}: {l_en}")
                c.setFillColor(muted); c.setFont(reg, 9)
                c.drawString(lx + 12, y + 12, l_tr)
                c.setFillColor(HexColor("#1C2230")); c.setStrokeColor(brand)
                c.roundRect(rx, y, cw, rh, 8, stroke=1, fill=1)
                c.setFillColor(text); c.setFont(bold, 11)
                c.drawString(rx + 12, y + rh - 22, f"{le}: {r_en}")
                c.setFillColor(brand); c.setFont(reg, 9)
                c.drawString(rx + 12, y + 12, r_tr)
                y -= rh + 12
        elif kind == "demo":
            c.setFillColor(panel); c.setStrokeColor(brand)
            c.roundRect(160, 120, 640, 230, 12, stroke=1, fill=1)
            c.setFillColor(brand); c.setFont(bold, 22)
            c.drawCentredString(W / 2, 300, "▶  LIVE / RECORDED DEMO")
            c.setFillColor(text); c.setFont(reg, 12)
            para(200, 265, s["en_l"], reg, 12, text, 560, 18)
            c.setFillColor(muted); c.setFont(reg, 10)
            para(200, 200, s["tr_l"], reg, 10, muted, 560, 15)
        elif kind == "closing":
            yy = para(40, H - 210, s["en_l"], reg, 15, text, W - 80, 22)
            para(40, yy - 6, s["tr_l"], reg, 12, muted, W - 80, 18)
            perimeter(70)
        c.showPage()

    c.save()
    print("wrote", PDF_OUT)


if __name__ == "__main__":
    build_pptx()
    build_pdf()
